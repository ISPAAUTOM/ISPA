import streamlit as st
import os
import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.dml import MSO_COLOR_TYPE, MSO_THEME_COLOR
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Cm, Pt
from docx import Document
from docx.shared import Pt as DocxPt, RGBColor as DocxRGBColor
import tempfile

# ============================================================================
# CONFIGURATION STREAMLIT
# ============================================================================

st.set_page_config(
    page_title="Modificateur de documents ISPA",
    page_icon="📄",
    layout="wide"
)

# ============================================================================
# CONFIGURATION DES STYLES ISPA
# ============================================================================

# Configuration du logo principal
LOGO_X = Cm(0.85)
LOGO_Y = Cm(0.98)
LOGO_WIDTH = Cm(2.73)
LOGO_HEIGHT = Cm(2.74)

# Configuration du favicon
FAVICON_X = Cm(4.29)
FAVICON_Y = Cm(4.14)
FAVICON_WIDTH = Cm(8.62)
FAVICON_HEIGHT = Cm(4.48)

# ---------------------------------------------------------------------------
# STYLES DE TEXTE (POWERPOINT)
# ---------------------------------------------------------------------------

TITRE_POLICE = "Lexend Bold"
TITRE_TAILLE = Pt(40)  # Changé de 42 à 40 selon vos specs
TITRE_COULEUR = RGBColor(111, 156, 235)  # Bleu #6F9CEB

CORPS_POLICE = "Lexend Regular"
CORPS_TAILLE = Pt(22)
CORPS_COULEUR = RGBColor(0, 0, 0)       # Noir #000000

# À partir du 2ème niveau (bullets)
BULLET_POLICE = "Lexend Light"
BULLET_TAILLE = Pt(18)
BULLET_COULEUR = RGBColor(0, 0, 0)      # Noir #000000

# ---------------------------------------------------------------------------
# SEUILS DE DÉTECTION LOGO
# ---------------------------------------------------------------------------

MAX_LEFT_LOGO = Cm(2)
MAX_TOP_LOGO = Cm(2)
MIN_RIGHT_FAVICON = Cm(43)
MAX_RIGHT_FAVICON = Cm(46)
MIN_TOP_FAVICON = Cm(3)
MAX_TOP_FAVICON = Cm(5)

# Doubles dimensions pour détection
DOUBLE_LOGO_WIDTH = 2 * LOGO_WIDTH
DOUBLE_LOGO_HEIGHT = 2 * LOGO_HEIGHT
DOUBLE_FAVICON_WIDTH = 2 * FAVICON_WIDTH
DOUBLE_FAVICON_HEIGHT = 2 * FAVICON_HEIGHT

# ========== STYLES WORD SPÉCIFIQUES ==========

WORD_TITLE_STYLE_NAMES = ["Title", "Titre 1", "Heading 1"]
WORD_SUBTIT_STYLE_NAMES = ["Subtitle", "Titre 2", "Heading 2"]

WORD_TITRE_POLICE = "Lexend Bold"
WORD_TITRE_TAILLE = DocxPt(28)
WORD_TITRE_COULEUR = DocxRGBColor(111, 156, 235)  # #6F9CEB

WORD_SOUS_TITRE_POLICE = "Lexend Light"
WORD_SOUS_TITRE_TAILLE = DocxPt(14)

WORD_TEXTE_POLICE = "Lexend Regular"
WORD_TEXTE_TAILLE = DocxPt(11)

# ---------------------------------------------------------------------------
# OUTILS PPTX
# ---------------------------------------------------------------------------

def remove_old_logo_if_small_in_corner(shape, progress_text):
    """Détecte et supprime l'ancien logo principal (coin haut-gauche)."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if (shape.left < MAX_LEFT_LOGO and shape.top < MAX_TOP_LOGO
                and shape.width <= DOUBLE_LOGO_WIDTH
                and shape.height <= DOUBLE_LOGO_HEIGHT):
                shape._element.getparent().remove(shape._element)
                progress_text.text(f"  → Ancien logo principal supprimé")
                return True
    except Exception as e:
        progress_text.text(f"  ⚠️ Erreur suppression logo: {str(e)}")
    return False

def remove_old_favicon_if_in_corner(shape, progress_text):
    """Détecte et supprime l'ancien favicon (coin haut-droit)."""
    try:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            if (MIN_RIGHT_FAVICON < shape.left < MAX_RIGHT_FAVICON 
                and MIN_TOP_FAVICON < shape.top < MAX_TOP_FAVICON
                and shape.width <= DOUBLE_FAVICON_WIDTH
                and shape.height <= DOUBLE_FAVICON_HEIGHT):
                shape._element.getparent().remove(shape._element)
                progress_text.text(f"  → Ancien favicon supprimé")
                return True
    except Exception as e:
        progress_text.text(f"  ⚠️ Erreur suppression favicon: {str(e)}")
    return False

def apply_paragraph_style(paragraph, font_name, font_size, font_color):
    try:
        if paragraph.font:
            paragraph.font.name = font_name
            paragraph.font.size = font_size
            if paragraph.font.color:
                paragraph.font.color.rgb = font_color

        for run in paragraph.runs:
            if run.font:
                run.font.name = font_name
                run.font.size = font_size
                if run.font.color:
                    run.font.color.rgb = font_color
    except Exception:
        pass

def appliquer_style_texte_pptx(text_frame, shape=None, progress_text=None):
    """Applique les styles ISPA aux textes."""
    if not text_frame:
        return

    try:
        text_frame.auto_size = MSO_AUTO_SIZE.NONE
    except:
        pass

    force_title = getattr(shape, "_force_title", False) if shape else False
    
    for paragraph in text_frame.paragraphs:
        try:
            if paragraph.level >= 1:
                # Bullet (2ème niveau et plus) => Lexend Light 18pt noir
                apply_paragraph_style(paragraph, BULLET_POLICE, BULLET_TAILLE, BULLET_COULEUR)
                if progress_text:
                    progress_text.text(f"    → [Bullet] {paragraph.text[:40]}...")
            else:
                if force_title:
                    # Titre => Lexend Bold 40pt bleu
                    apply_paragraph_style(paragraph, TITRE_POLICE, TITRE_TAILLE, TITRE_COULEUR)
                    if progress_text:
                        progress_text.text(f"    → [Titre 40pt] {paragraph.text[:40]}...")
                else:
                    # Corps => Lexend Regular 22pt noir
                    apply_paragraph_style(paragraph, CORPS_POLICE, CORPS_TAILLE, CORPS_COULEUR)
                    if progress_text:
                        progress_text.text(f"    → [Corps 22pt] {paragraph.text[:40]}...")
        except Exception:
            pass

def style_table(table):
    try:
        for row in table.rows:
            for cell in row.cells:
                if cell.text_frame:
                    for paragraph in cell.text_frame.paragraphs:
                        apply_paragraph_style(paragraph, CORPS_POLICE, CORPS_TAILLE, CORPS_COULEUR)
    except:
        pass

def get_text_content(shape):
    """Retourne le texte d'un shape (strip) ou ""."""
    try:
        if hasattr(shape, "text_frame") and shape.text_frame:
            return shape.text_frame.text.strip()
    except:
        pass
    return ""

def traiter_pptx(fichier_entree, logo_path, favicon_path, progress_bar, progress_text):
    """Traite un fichier PowerPoint avec logo et favicon."""
    try:
        progress_text.text("Ouverture du fichier PowerPoint...")
        pres = Presentation(fichier_entree)
        
        total_slides = len(pres.slides)
        master_logo_removed = False
        master_favicon_removed = False

        # MASTER
        progress_text.text("Traitement des masters...")
        for master_idx, master in enumerate(pres.slide_masters, start=1):
            progress_text.text(f"Master {master_idx}")
            
            any_logo_removed = False
            any_favicon_removed = False
            
            for shape in list(master.shapes):
                if remove_old_logo_if_small_in_corner(shape, progress_text):
                    any_logo_removed = True
                elif remove_old_favicon_if_in_corner(shape, progress_text):
                    any_favicon_removed = True

            if any_logo_removed:
                master_logo_removed = True
            if any_favicon_removed:
                master_favicon_removed = True

            # Styles texte sur master
            text_shapes = [sh for sh in master.shapes if hasattr(sh, "text_frame")]
            text_shapes.sort(key=lambda s: s.top)
            filtered = []
            for sh in text_shapes:
                txt = get_text_content(sh)
                if len(txt) > 3:
                    filtered.append(sh)

            if filtered:
                filtered[0]._force_title = True

            for sh in text_shapes:
                appliquer_style_texte_pptx(sh.text_frame, sh, progress_text)

        # SLIDES
        slides_list = list(pres.slides)
        for idx, slide in enumerate(slides_list, start=1):
            progress_bar.progress(idx / total_slides)
            progress_text.text(f"Slide {idx}/{total_slides}")

            old_logo_removed = False
            old_favicon_removed = False
            
            for shape in list(slide.shapes):
                if remove_old_logo_if_small_in_corner(shape, progress_text):
                    old_logo_removed = True
                elif remove_old_favicon_if_in_corner(shape, progress_text):
                    old_favicon_removed = True

            # Insérer nouveau logo si nécessaire
            if old_logo_removed or master_logo_removed:
                try:
                    slide.shapes.add_picture(logo_path, LOGO_X, LOGO_Y, 
                                           width=LOGO_WIDTH, height=LOGO_HEIGHT)
                    progress_text.text("  → Nouveau logo inséré")
                except Exception as e:
                    progress_text.text(f"  ⚠️ Erreur insertion logo: {str(e)}")

            # Insérer nouveau favicon si nécessaire
            if old_favicon_removed or master_favicon_removed:
                if favicon_path:  # Seulement si un favicon a été fourni
                    try:
                        slide.shapes.add_picture(favicon_path, FAVICON_X, FAVICON_Y, 
                                               width=FAVICON_WIDTH, height=FAVICON_HEIGHT)
                        progress_text.text("  → Nouveau favicon inséré")
                    except Exception as e:
                        progress_text.text(f"  ⚠️ Erreur insertion favicon: {str(e)}")

            # Traitement des textes
            text_shapes = []
            for sh in slide.shapes:
                if hasattr(sh, "text_frame"):
                    text_shapes.append(sh)
                if sh.shape_type == MSO_SHAPE_TYPE.TABLE:
                    style_table(sh.table)

            text_shapes.sort(key=lambda s: s.top)
            filtered = []
            for sh in text_shapes:
                txt = get_text_content(sh)
                if len(txt) > 3:
                    filtered.append(sh)

            if filtered:
                filtered[0]._force_title = True

            for sh in text_shapes:
                appliquer_style_texte_pptx(sh.text_frame, sh, progress_text)

        # Sauvegarder
        output = io.BytesIO()
        pres.save(output)
        output.seek(0)
        
        progress_bar.progress(1.0)
        progress_text.text("✅ PowerPoint traité avec succès!")
        
        return output

    except Exception as e:
        st.error(f"❌ Erreur PowerPoint: {str(e)}")
        return None

# ============================================================================
# TRAITEMENT WORD
# ============================================================================

def apply_run_style_word(run, font_name, font_size, font_color=None):
    try:
        if run.font:
            run.font.name = font_name
            run.font.size = font_size
            if font_color:
                run.font.color.rgb = font_color
            else:
                run.font.color.rgb = DocxRGBColor(0, 0, 0)
    except:
        pass

def style_word_paragraph_by_name(paragraph):
    try:
        if paragraph.style and paragraph.style.name:
            style_name = paragraph.style.name.lower()
            if any(s.lower() in style_name for s in WORD_TITLE_STYLE_NAMES):
                for run in paragraph.runs:
                    apply_run_style_word(run, WORD_TITRE_POLICE, WORD_TITRE_TAILLE, WORD_TITRE_COULEUR)
                return "TITLE"
            elif any(s.lower() in style_name for s in WORD_SUBTIT_STYLE_NAMES):
                for run in paragraph.runs:
                    apply_run_style_word(run, WORD_SOUS_TITRE_POLICE, WORD_SOUS_TITRE_TAILLE)
                return "SUB"
    except:
        pass
    return None

def appliquer_style_texte_word(paragraph, is_title_fallback=False):
    txt = paragraph.text.strip()
    if not txt:
        return

    style_result = style_word_paragraph_by_name(paragraph)
    if style_result == "TITLE" or style_result == "SUB":
        return

    is_bullet = (txt.startswith("- ") or txt.startswith("* "))
    if is_bullet:
        for run in paragraph.runs:
            apply_run_style_word(run, WORD_SOUS_TITRE_POLICE, WORD_SOUS_TITRE_TAILLE)
    else:
        if is_title_fallback:
            for run in paragraph.runs:
                apply_run_style_word(run, WORD_TITRE_POLICE, WORD_TITRE_TAILLE, WORD_TITRE_COULEUR)
        else:
            for run in paragraph.runs:
                apply_run_style_word(run, WORD_TEXTE_POLICE, WORD_TEXTE_TAILLE)

def traiter_docx(fichier_entree, logo_path, favicon_path, progress_bar, progress_text):
    try:
        progress_text.text("Ouverture du fichier Word...")
        doc = Document(fichier_entree)

        # En-têtes - on ne gère que le logo principal pour Word
        progress_text.text("Traitement des en-têtes...")
        for section_idx, section in enumerate(doc.sections, start=1):
            progress_bar.progress(0.3)
            header = section.header
            logo_found = False

            for para in header.paragraphs:
                for run in para.runs:
                    if run._element.findall('.//w:drawing', namespaces=run._element.nsmap):
                        logo_found = True
                        run._element.clear()

            if logo_found:
                run = header.paragraphs[0].add_run()
                run.add_picture(logo_path, width=LOGO_WIDTH, height=LOGO_HEIGHT)
                progress_text.text("✅ Nouveau logo ajouté à l'en-tête")

        # 1ère image corps
        progress_text.text("Recherche et remplacement du logo...")
        progress_bar.progress(0.5)
        found_first_image = False
        for para in doc.paragraphs:
            if not found_first_image:
                for run in para.runs:
                    if run._element.findall('.//w:drawing', namespaces=run._element.nsmap):
                        run._element.clear()
                        new_run = para.add_run()
                        new_run.add_picture(logo_path, width=LOGO_WIDTH, height=LOGO_HEIGHT)
                        progress_text.text("✅ Logo remplacé dans le document")
                        found_first_image = True
                        break
            else:
                break

        # Styles
        progress_text.text("Application des styles...")
        progress_bar.progress(0.7)
        found_title = False
        for para in doc.paragraphs:
            txt = para.text.strip()
            if not txt:
                continue
            if not found_title:
                appliquer_style_texte_word(para, is_title_fallback=True)
                found_title = True
            else:
                appliquer_style_texte_word(para, is_title_fallback=False)

        # Sauvegarder
        output = io.BytesIO()
        doc.save(output)
        output.seek(0)
        
        progress_bar.progress(1.0)
        progress_text.text("✅ Document Word traité avec succès!")
        
        return output

    except Exception as e:
        st.error(f"❌ Erreur Word: {str(e)}")
        return None

# ============================================================================
# INTERFACE STREAMLIT
# ============================================================================

def main():
    st.title("🎨 Modificateur de documents ISPA")
    st.markdown("### Transformez vos présentations PowerPoint et documents Word")
    
    # Info box pour les spécificités ISPA
    with st.info("ℹ️ Spécificités ISPA"):
        st.markdown("""
        - **Titres** : Lexend Bold 40pt - Couleur #6F9CEB
        - **Corps** : Lexend Regular 22pt - Noir
        - **Bullets** : Lexend Light 18pt - Noir
        - **Logo** : Coin supérieur gauche
        - **Favicon** : Coin supérieur droit (PowerPoint uniquement)
        """)
    
    col1, col2, col3 = st.columns(3)
    
    with col1:
        st.markdown("#### 📁 Document à traiter")
        uploaded_file = st.file_uploader(
            "Glissez votre fichier ici",
            type=['pptx', 'docx'],
            help="Formats supportés : PowerPoint (.pptx) et Word (.docx)"
        )
    
    with col2:
        st.markdown("#### 🖼️ Logo principal")
        logo_file = st.file_uploader(
            "Logo (obligatoire)",
            type=['png', 'jpg', 'jpeg'],
            help="Sera placé en haut à gauche"
        )
    
    with col3:
        st.markdown("#### 🔰 Favicon")
        favicon_file = st.file_uploader(
            "Favicon (optionnel)",
            type=['png', 'jpg', 'jpeg'],
            help="PowerPoint uniquement - haut à droite"
        )
    
    if uploaded_file and logo_file:
        st.markdown("---")
        
        # Afficher les infos du fichier
        file_details = {
            "Nom du fichier": uploaded_file.name,
            "Type": uploaded_file.type,
            "Taille": f"{uploaded_file.size / 1024:.1f} KB"
        }
        st.json(file_details)
        
        # Avertissement si favicon pour Word
        if uploaded_file.name.lower().endswith('.docx') and favicon_file:
            st.warning("⚠️ Le favicon n'est pas supporté pour les documents Word")
        
        if st.button("🚀 Lancer le traitement", type="primary"):
            
            # Créer des fichiers temporaires
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_logo:
                tmp_logo.write(logo_file.getbuffer())
                logo_path = tmp_logo.name
            
            favicon_path = None
            if favicon_file and uploaded_file.name.lower().endswith('.pptx'):
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as tmp_favicon:
                    tmp_favicon.write(favicon_file.getbuffer())
                    favicon_path = tmp_favicon.name
            
            # Progress bars
            progress_bar = st.progress(0)
            progress_text = st.empty()
            
            # Traitement selon le type
            if uploaded_file.name.lower().endswith('.pptx'):
                with st.spinner('Traitement du PowerPoint en cours...'):
                    output = traiter_pptx(uploaded_file, logo_path, favicon_path, progress_bar, progress_text)
            else:
                with st.spinner('Traitement du Word en cours...'):
                    output = traiter_docx(uploaded_file, logo_path, favicon_path, progress_bar, progress_text)
            
            # Nettoyer les fichiers temporaires
            os.unlink(logo_path)
            if favicon_path:
                os.unlink(favicon_path)
            
            if output:
                st.success("✅ Traitement terminé avec succès!")
                
                # Bouton de téléchargement
                st.download_button(
                    label="📥 Télécharger le fichier modifié",
                    data=output,
                    file_name=f"ISPA_{uploaded_file.name}",
                    mime="application/octet-stream",
                    type="primary"
                )
    
    # Footer
    st.markdown("---")
    st.markdown("""
    <div style='text-align: center; color: #666;'>
        <small>
        💡 Le script détecte et remplace automatiquement les anciens logos et favicons,
        puis applique la charte graphique ISPA.
        </small>
    </div>
    """, unsafe_allow_html=True)

if __name__ == "__main__":
    main()
